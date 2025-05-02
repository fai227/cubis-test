from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import re

index = 1
def to_cpp_variable_name(input_str):
    if input_str == "#":
        return "CHANGABLE_TEXT"
    
    global index

    # + -> plus, - -> minus, / -> or, & -> and
    input_str = input_str.replace('+', 'plus')
    input_str = input_str.replace('-', 'minus')
    input_str = input_str.replace('/', 'or')
    input_str = input_str.replace('&', 'and')

    # スペースをアンダースコアに変換
    valid_name = input_str.replace(' ', '_')
    
    # 数字とアルファベット以外を削除
    valid_name = re.sub(r'[^a-zA-Z0-9_]', '', valid_name)
    
    # 連続するアンダースコアを1つに変換
    valid_name = re.sub(r'_{2,}', '_', valid_name)
    
    # アルファベットを大文字に変換
    valid_name = valid_name.upper()
    
    # 変数名が数字で始まる場合、先頭にアンダースコアを追加
    if valid_name and valid_name[0].isdigit():
        valid_name = '_' + valid_name
    
    # アンダースコアのみ、数字のみ、またはアンダースコアと数字の組み合わせの場合の処理
    if not valid_name or re.fullmatch(r'[_0-9]+', valid_name):
        valid_name = f"NO_TEXT_{index}"
        index += 1

    return valid_name


def extract_ui_elements(pptx_path):
    prs = Presentation(pptx_path)
    slide_width = prs.slide_width.pt
    slide_height = prs.slide_height.pt

    # スライドごとのUI要素
    all_ui_elements = []  
    constant_texts = {}  # 定数テキスト    

    # すべてのスライドを処理
    for slide_index, slide in enumerate(prs.slides):
        slide_ui_elements = []  # このスライドのUI要素

        text_count = 1
        button_count = 1

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            if shape.shape_type not in [MSO_SHAPE_TYPE.TEXT_BOX, MSO_SHAPE_TYPE.AUTO_SHAPE]:
                continue

            text = shape.text.strip()
            if not text:
                continue

            element_type = "text" if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX else "button"

            left = shape.left.pt
            top = shape.top.pt
            width = shape.width.pt
            height = shape.height.pt

            # 相対中心座標 (0.0〜1.0)
            center_x = (left + width / 2) / slide_width
            center_y = (top + height / 2) / slide_height
            rel_width = width / slide_width
            rel_height = height / slide_height

            if element_type == "text":
                name = element_type + str(text_count)
                text_count += 1
                text_var_name = to_cpp_variable_name(text)
                if not text in constant_texts:
                    constant_texts[text] = text_var_name
                else:
                    text_var_name = constant_texts[text]
                slide_ui_elements.append({
                    "type": element_type,
                    "name": name,
                    "center_x": center_x,
                    "center_y": center_y,
                    "rel_width": rel_width,
                    "rel_height": rel_height,
                    "text_var_name": text_var_name
                })
            else:
                name = element_type + str(button_count)
                button_count += 1
                button_var_name = to_cpp_variable_name(text)
                if not text in constant_texts:
                    constant_texts[text] = button_var_name
                else:
                    button_var_name = constant_texts[text]
                slide_ui_elements.append({
                    "type": element_type,
                    "name": name,
                    "center_x": center_x,
                    "center_y": center_y,
                    "rel_width": rel_width,
                    "rel_height": rel_height,
                    "text_var_name": button_var_name
                })
                

        # スライドごとのUI要素を追加
        all_ui_elements.append(slide_ui_elements)

    return all_ui_elements, constant_texts

import math

def search_best_ui_element(elements, element, direction="right", angle_limit_deg=45):
    angle_limit_cos = math.cos(math.radians(angle_limit_deg))

    # 基準方向ベクトル
    if direction == "right":
        ref_dir = (1, 0)
    elif direction == "left":
        ref_dir = (-1, 0)
    elif direction == "up":
        ref_dir = (0, -1)
    elif direction == "down":
        ref_dir = (0, 1)
    else:
        raise ValueError("Unsupported direction")

    best_element = None
    best_distance_sq = float("inf")

    for e in elements:
        if e == element:
            continue
        if e["type"] != element["type"]:
            continue

        dx = e["center_x"] - element["center_x"]
        dy = e["center_y"] - element["center_y"]

        len_sq = dx * dx + dy * dy
        if len_sq == 0:
            continue

        # 正規化ベクトルと内積で角度判定
        nx, ny = dx / math.sqrt(len_sq), dy / math.sqrt(len_sq)
        dot = nx * ref_dir[0] + ny * ref_dir[1]

        if dot > angle_limit_cos:
            if len_sq < best_distance_sq:
                best_distance_sq = len_sq
                best_element = e
    if best_element is None:
        return "nullptr"
    return best_element["name"]


def solve_connection(all_ui_elements):
    for ui_elements in all_ui_elements:
        for ui_element in ui_elements:
            if ui_element["type"] == "text":
                continue
            ui_element["connection"] = {
                "up": search_best_ui_element(ui_elements, ui_element, direction="up"),
                "down": search_best_ui_element(ui_elements, ui_element, direction="down"),
                "left": search_best_ui_element(ui_elements, ui_element, direction="left"),
                "right": search_best_ui_element(ui_elements, ui_element, direction="right"),
            }

def export_cpp_code(all_ui_elements, constant_texts, output_path):
    with open(output_path, "w", encoding="utf-8") as out:
        # スライドごとにMarkdown形式でC++コードを出力
        for slide_index, ui_elements in enumerate(all_ui_elements):
            # 進捗表示
            out.write(f"# Slide {slide_index + 1}\n\n")
            out.write("```c++\n")

            # 各UI要素を処理
            for ui_element in ui_elements:
                elem_type = ui_element["type"]
                name = ui_element["name"]
                center_x = ui_element["center_x"]
                center_y = ui_element["center_y"]
                rel_width = ui_element["rel_width"]
                rel_height = ui_element["rel_height"]
                var_name = ui_element["text_var_name"]

                if elem_type == "button":
                    out.write(f"Button *{name} = new Button({{{center_x:.3f}, {center_y:.3f}, {rel_width:.3f}, {rel_height:.3f}}}, UIElementDisplayText::{var_name}.c_str());\n")
                else:
                    out.write(f"Text *{name} = new Text({{{center_x:.3f}, {center_y:.3f}, {rel_width:.3f}, {rel_height:.3f}}}, UIElementDisplayText::{var_name}.c_str());\n")
            out.write("\n")

            # 各UIのナビゲーションの出力
            for ui_element in ui_elements:
                name = ui_element["name"]
                connection = ui_element.get("connection")

                if connection is None:
                    continue
                out.write(f"{name}->SetNavigation(")

                # 各方向の接続先を出力
                tmp_list = []
                for direction in ["up", "down", "left", "right"]:
                    target = connection[direction]
                    tmp_list.append(target)
                out.write(", ".join(tmp_list))
                out.write(");\n")

            # uiElements.push_back() の出力
            out.write("\n")
            for ui_element in ui_elements:
                name = ui_element["name"]
                out.write(f"uiElements.push_back({name});\n")

            out.write("```\n\n")

        # Textsセクションの出力
        out.write("# Texts\n\n")
        out.write("```c++\n")
        out.write("namespace UIElementDisplayText\n{\n")
        for text,var_name  in constant_texts.items():
            out.write(f'    const std::string {var_name} = "{text}";\n')
        out.write("}\n")
        out.write("```\n")


if __name__ == "__main__":
    pptx_path = "blueprint.pptx"
    output_path = "cpp.md"
    all_ui_elements, constant_texts = extract_ui_elements(pptx_path)
    solve_connection(all_ui_elements)
    export_cpp_code(all_ui_elements, constant_texts, output_path)